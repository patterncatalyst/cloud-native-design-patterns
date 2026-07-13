#!/usr/bin/env python3
"""Add Appendix P slides to the .NET PPTX presentation."""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

PPTX = "../lgtm-presentation/Designing-Cloud-Native-APIs-DotNet.pptx"

prs = Presentation(PPTX)
layout = prs.slide_layouts[0]

# Extract background image and logos from Appendix J divider (slide 206)
ref_divider = prs.slides[205]
ref_content = prs.slides[206]  # slide 207

bg_blob = None
bg_ct = None
div_logo_blob = None
div_logo_ct = None
content_logo_blob = None
content_logo_ct = None

for shape in ref_divider.shapes:
    if hasattr(shape, 'image'):
        if shape.width > 12000000:
            bg_blob = shape.image.blob
            bg_ct = shape.image.content_type
        else:
            div_logo_blob = shape.image.blob
            div_logo_ct = shape.image.content_type

for shape in ref_content.shapes:
    if hasattr(shape, 'image') and shape.width < 2000000:
        content_logo_blob = shape.image.blob
        content_logo_ct = shape.image.content_type
        break

last_page = 256


def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                bold=False, color=RGBColor(0x15, 0x15, 0x15), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_name, font_size, bold=False,
                          color=RGBColor(0x15, 0x15, 0x15), line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.space_after = Pt(6)
        if p.runs:
            run = p.runs[0]
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox


def add_content_slide(section_tag, title, body_lines, footer, page_num):
    slide = prs.slides.add_slide(layout)

    # Section header (red, top)
    add_textbox(slide, 566928, 384048, 11057700, 292500,
                section_tag,
                "Red Hat Text", Pt(12), bold=True, color=RGBColor(0xEE, 0x00, 0x00))

    # Title
    add_textbox(slide, 566928, 676656, 11057700, 822900,
                title,
                "Overpass SemiBold", Pt(30), bold=True, color=RGBColor(0x15, 0x15, 0x15))

    # Body
    add_multiline_textbox(slide, 566928, 1600200, 11057700, 4480500,
                          body_lines,
                          "Red Hat Text", Pt(17), bold=False,
                          color=RGBColor(0x15, 0x15, 0x15))

    # Footer
    add_textbox(slide, 566928, 5943600, 11057700, 310800,
                footer,
                "Red Hat Text", Pt(12), bold=False, color=RGBColor(0x5A, 0x5A, 0x5A))

    # Page number
    add_textbox(slide, 566928, 6364224, 914400, 274200,
                str(page_num),
                "Red Hat Text", Pt(10), bold=False, color=RGBColor(0x8A, 0x8A, 0x8A))

    # Logo
    if content_logo_blob:
        from io import BytesIO
        slide.shapes.add_picture(BytesIO(content_logo_blob),
                                 Emu(10442448), Emu(6291072),
                                 Emu(1216152), Emu(283464))

    return slide


def add_diagram_slide(section_tag, title, svg_path, footer, page_num):
    slide = prs.slides.add_slide(layout)

    # Section header
    add_textbox(slide, 566928, 384048, 11057700, 292500,
                section_tag,
                "Red Hat Text", Pt(12), bold=True, color=RGBColor(0xEE, 0x00, 0x00))

    # Title
    add_textbox(slide, 566928, 676656, 11057700, 822900,
                title,
                "Overpass SemiBold", Pt(30), bold=True, color=RGBColor(0x15, 0x15, 0x15))

    # Diagram image (centered, large)
    from io import BytesIO
    import subprocess, os
    # Convert SVG to PNG for PPTX embedding
    png_path = svg_path.replace('.svg', '.png')
    if not os.path.exists(png_path):
        try:
            subprocess.run(['rsvg-convert', '-w', '1600', svg_path, '-o', png_path],
                           check=True, capture_output=True)
        except FileNotFoundError:
            # rsvg-convert not available, try cairosvg
            try:
                import cairosvg
                cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=1600)
            except ImportError:
                print(f"  WARNING: Cannot convert {svg_path} to PNG - skipping image")
                png_path = None

    if png_path and os.path.exists(png_path):
        from PIL import Image
        img = Image.open(png_path)
        img_w, img_h = img.size
        # Scale to fit in the content area (roughly 8906256 x 4343400 EMU)
        max_w = 8906256
        max_h = 4343400
        scale = min(max_w / (img_w * 914400 / 96), max_h / (img_h * 914400 / 96))
        pic_w = int(img_w * 914400 / 96 * scale)
        pic_h = int(img_h * 914400 / 96 * scale)
        pic_left = (12191675 - pic_w) // 2
        pic_top = 1600200
        slide.shapes.add_picture(png_path, Emu(pic_left), Emu(pic_top),
                                 Emu(pic_w), Emu(pic_h))

    # Footer
    add_textbox(slide, 566928, 5943600, 11057700, 310800,
                footer,
                "Red Hat Text", Pt(12), bold=False, color=RGBColor(0x5A, 0x5A, 0x5A))

    # Page number
    add_textbox(slide, 566928, 6364224, 914400, 274200,
                str(page_num),
                "Red Hat Text", Pt(10), bold=False, color=RGBColor(0x8A, 0x8A, 0x8A))

    # Logo
    if content_logo_blob:
        slide.shapes.add_picture(BytesIO(content_logo_blob),
                                 Emu(10442448), Emu(6291072),
                                 Emu(1216152), Emu(283464))

    return slide


# ---- Section divider slide ----
from io import BytesIO

divider = prs.slides.add_slide(layout)

# Background image
if bg_blob:
    divider.shapes.add_picture(BytesIO(bg_blob),
                               Emu(0), Emu(0),
                               Emu(12191694), Emu(6857999))

# Letter "P"
add_textbox(divider, 5733288, 2121408, 5852100, 457200,
            "P",
            "Overpass SemiBold", Pt(22), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# Title
add_textbox(divider, 5705856, 2596896, 6126600, 1371600,
            "Appendix P · .NET on Linux",
            "Overpass SemiBold", Pt(40), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# Subtitle
add_textbox(divider, 5733288, 4114800, 6035100, 548700,
            "Runtime, Kestrel, and Containers — what changes when Windows goes away.",
            "Red Hat Text", Pt(15), bold=False, color=RGBColor(0xFF, 0xD9, 0xD9))

# Divider logo
if div_logo_blob:
    divider.shapes.add_picture(BytesIO(div_logo_blob),
                               Emu(10442448), Emu(6291072),
                               Emu(1216152), Emu(283464))

page = last_page + 1
DIAGRAMS = "../assets/diagrams"

# ---- Slide 1: C# vs .NET ----
add_content_slide(
    "APPENDIX P · RUNTIME",
    "C# is a language. .NET is a platform.",
    [
        "C# is a statically-typed, garbage-collected language — one of several that target .NET. F# and VB.NET are the others. All three compile to the same intermediate representation.",
        "",
        ".NET (since .NET 5) is the platform: CoreCLR runtime, Base Class Library (System.*, Microsoft.Extensions.*), the Roslyn compiler, NuGet, and the dotnet CLI. The platform is cross-platform.",
        "",
        "When this book says “.NET 10,” it means the platform. When it says “C#,” it means the language. A .csproj targets a .NET version (net10.0), not a C# version, because the platform version determines what APIs are available.",
    ],
    "The distinction matters because conflating them causes architectural confusion — the runtime is language-agnostic.",
    page,
)
page += 1

# ---- Slide 2: CIL + compilation pipeline diagram ----
add_diagram_slide(
    "APPENDIX P · COMPILATION",
    "CIL: what dotnet build actually produces",
    f"{DIAGRAMS}/29-dotnet-compilation.svg",
    "Roslyn compiles C# to CIL bytecode — RyuJIT compiles CIL to native code method-at-a-time at runtime.",
    page,
)
page += 1

# ---- Slide 3: CoreCLR runtime ----
add_content_slide(
    "APPENDIX P · RUNTIME",
    "CoreCLR on Linux: libcoreclr.so",
    [
        "JIT compilation (RyuJIT) — CIL to native x86-64 or ARM64, one method at a time, first call only.",
        "Garbage collection — generational, concurrent GC. Server GC mode (default in ASP.NET Core) allocates one heap per CPU core.",
        "Type system — reflection, generics, interface dispatch.",
        "Threading — maps .NET threads to pthreads; Task and async/await run on the thread pool.",
        "P/Invoke — call native C libraries directly (libpq, librdkafka, openssl).",
        "",
        "NativeAOT alternative: ahead-of-time compilation to a single native binary. Faster startup, less memory, but no Reflection.Emit — JIT is the default for full library compatibility.",
    ],
    "After JIT, a .NET method is indistinguishable from one compiled by GCC. The runtime stays for GC and exceptions, but the hot path is native code.",
    page,
)
page += 1

# ---- Slide 4: Kestrel pipeline diagram ----
add_diagram_slide(
    "APPENDIX P · KESTREL",
    "Kestrel: the web server that replaced IIS",
    f"{DIAGRAMS}/29-kestrel-pipeline.svg",
    "Built into ASP.NET Core — no IIS, no nginx needed. Listens on TCP, handles TLS via OpenSSL, parses HTTP, runs the middleware pipeline.",
    page,
)
page += 1

# ---- Slide 5: Kestrel details ----
add_content_slide(
    "APPENDIX P · KESTREL",
    "Why not nginx or Apache in front?",
    [
        "In Kubernetes, the service mesh handles TLS termination, load balancing, and retries. Kestrel handles everything else.",
        "",
        "Adding nginx as a sidecar adds latency, memory, and operational complexity for features Kestrel already provides.",
        "",
        "Every example in this book: builder.WebHost.UseUrls(\"http://+:8080\")",
        "For gRPC (HTTP/2 only): builder.WebHost.ConfigureKestrel(k => k.ListenAnyIP(50051, o => o.Protocols = HttpProtocols.Http2))",
        "",
        "Port 8080 because ports below 1024 require root, which containers should not run as.",
    ],
    "Kestrel supports HTTP/1.1, HTTP/2, HTTP/3 natively — gRPC, WebSockets, and REST share the same host.",
    page,
)
page += 1

# ---- Slide 6: Container layers diagram ----
add_diagram_slide(
    "APPENDIX P · CONTAINERS",
    "Multi-stage container build",
    f"{DIAGRAMS}/29-container-layers.svg",
    "SDK compiles (∼780 MB image), only the runtime + app DLLs ship (∼200 MB image). The SDK never reaches production.",
    page,
)
page += 1

# ---- Slide 7: UBI and what ships ----
add_content_slide(
    "APPENDIX P · CONTAINERS",
    "Red Hat UBI: what ships to production",
    [
        "UBI9 base (∼70 MB) — glibc, openssl, ca-certificates, tzdata, libicu.",
        "ASP.NET Core runtime (∼100 MB) — libcoreclr.so, libclrjit.so, BCL assemblies, ASP.NET shared framework.",
        "Your application (∼5–80 MB) — compiled .dll files, NuGet package DLLs, appsettings.json.",
        "",
        "Not in the image: SDK, MSBuild, Roslyn, NuGet, source code, obj/ and bin/ intermediates, test projects.",
        "",
        "UBI images are freely redistributable, based on RHEL 9, and supported by Red Hat on OpenShift. They contain the same .NET runtime bits as Microsoft’s images — the difference is the base OS layer.",
    ],
    "ubi9/dotnet-100 (SDK, build only) → ubi9/dotnet-100-aspnet (runtime, ships) — both from registry.access.redhat.com.",
    page,
)
page += 1

# ---- Slide 8: Container runtime lifecycle ----
add_content_slide(
    "APPENDIX P · LIFECYCLE",
    "Container runtime lifecycle",
    [
        "1. CoreCLR initializes — loads libcoreclr.so, starts the GC, initializes the thread pool.",
        "2. Host builder runs — registers services (DI), configures logging, sets up OpenTelemetry.",
        "3. Kestrel binds to port 8080 — begins accepting HTTP connections.",
        "4. Application is ready — Kubernetes starts sending liveness and readiness probes.",
        "5. On shutdown, SIGTERM arrives — .NET host catches it and begins graceful shutdown: stops accepting new requests, drains in-flight, flushes telemetry, disposes services.",
        "",
        "The entire lifecycle is the same as any Linux process. No IIS app pool recycle, no w3wp.exe, no application domain.",
    ],
    "The container is the process boundary — SIGTERM is your shutdown signal, HostOptions.ShutdownTimeout is your budget.",
    page,
)
page += 1

# ---- Slide 9: Take-aways ----
add_content_slide(
    "APPENDIX P · TAKE-AWAYS",
    "Take-aways",
    [
        "C# compiles to CIL bytecode, not native code. The CoreCLR JIT (RyuJIT) compiles to native at runtime, method-at-a-time. After that first call, performance is identical to C or C++.",
        "Kestrel is the built-in web server — it handles TLS (via OpenSSL), HTTP/1.1+2+3, gRPC, and WebSockets natively. No reverse proxy layer is needed inside the Pod.",
        "Multi-stage Dockerfiles separate the SDK (∼780 MB) from the runtime (∼170 MB). Only the published DLLs and the ASP.NET runtime ship to production.",
        "Red Hat UBI images are the same .NET runtime on a RHEL 9 base — freely redistributable, commercially supported on OpenShift.",
    ],
    "References: .NET runtime architecture (docs.microsoft.com) · Kestrel web server (ASP.NET Core docs) · Red Hat UBI .NET images (access.redhat.com)",
    page,
)

prs.save(PPTX)
print(f"Saved {page - last_page} new slides (pages {last_page + 1}–{page}) to {PPTX}")
